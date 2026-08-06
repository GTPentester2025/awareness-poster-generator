from pathlib import Path

import pytest
from pptx import Presentation

from app.builder import A4_EMU, build
from tests.test_content import VALID

PORTRAIT = (7_560_000, 10_692_000)


def _texts(pptx_path):
    prs = Presentation(str(pptx_path))
    out = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            out.append(shape.text_frame.text)
    return prs, out


def test_a4_constants():
    assert A4_EMU["portrait"] == PORTRAIT
    assert A4_EMU["landscape"] == (PORTRAIT[1], PORTRAIT[0])


@pytest.mark.parametrize("orientation", ["portrait", "landscape"])
def test_build_no_image(tmp_path, orientation):
    p = build(VALID, None, orientation, tmp_path)
    prs, texts = _texts(p)
    w, h = A4_EMU[orientation]
    assert (prs.slide_width, prs.slide_height) == (w, h)
    joined = "\n".join(texts)
    assert VALID["headline"] in joined
    assert VALID["cta"] in joined
    for fact in VALID["facts"]:
        assert fact in joined


def test_build_with_image(tmp_path):
    import base64
    from tests.test_artwork import PNG_1PX
    img = tmp_path / "bg.png"
    img.write_bytes(base64.b64decode(PNG_1PX))
    p = build(VALID, img, "portrait", tmp_path)
    prs, _ = _texts(p)
    pics = [s for s in prs.slides[0].shapes if s.shape_type == 13]  # PICTURE
    assert len(pics) == 1
    assert (pics[0].width, pics[0].height) == A4_EMU["portrait"]


def test_bad_orientation(tmp_path):
    with pytest.raises(ValueError):
        build(VALID, None, "square", tmp_path)
