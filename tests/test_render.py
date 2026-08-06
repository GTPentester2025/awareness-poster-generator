import base64

import pytest
from pptx import Presentation

from app.builder import A4_EMU, render
from app.design import ARCHETYPES
from tests.test_design import VALID_SPEC
from tests.test_artwork import PNG_1PX

VARIANT = {
    "angle": "precautions",
    "headline": "Stay Safe on the Roads",
    "subheadline": "Small habits, big difference",
    "points": [
        {"stat": "100%", "text": "Buckle up every trip"},
        {"stat": "5x", "text": "Never text and drive"},
        {"stat": "-20", "text": "Slow down in the rain"},
        {"stat": "0", "text": "Zero drinks before driving"},
        {"stat": "2s", "text": "Keep a two-second gap"},
    ],
    "cta": "Drive aware. Arrive alive.",
    "sources": ["WHO Road Safety 2023"],
}


def _texts(path):
    prs = Presentation(str(path))
    slide = prs.slides[0]
    joined = "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    return prs, slide, joined


@pytest.mark.parametrize("archetype", sorted(ARCHETYPES))
@pytest.mark.parametrize("orientation", ["portrait", "landscape"])
def test_render_every_archetype_solid(tmp_path, archetype, orientation):
    spec = {**VALID_SPEC, "archetype": archetype, "background_style": "solid"}
    p = render(spec, VARIANT, None, orientation, tmp_path)
    prs, slide, joined = _texts(p)
    assert (prs.slide_width, prs.slide_height) == A4_EMU[orientation]
    assert VARIANT["headline"] in joined
    assert VARIANT["cta"] in joined
    for point in VARIANT["points"]:
        assert point["text"] in joined


@pytest.mark.parametrize("n_points", [3, 4, 6])
def test_render_variable_point_counts(tmp_path, n_points):
    variant = {**VARIANT, "points": VARIANT["points"][:n_points] if n_points <= 5
               else VARIANT["points"] + [{"stat": "+1", "text": "Extra safety habit"}]}
    spec = {**VALID_SPEC, "archetype": "hero_top", "background_style": "solid"}
    p = render(spec, variant, None, "portrait", tmp_path)
    _, _, joined = _texts(p)
    for point in variant["points"]:
        assert point["text"] in joined


def test_render_sources_footer(tmp_path):
    spec = {**VALID_SPEC, "background_style": "solid"}
    p = render(spec, VARIANT, None, "portrait", tmp_path)
    _, _, joined = _texts(p)
    assert "WHO Road Safety 2023" in joined


def test_render_bad_orientation(tmp_path):
    with pytest.raises(ValueError):
        render(VALID_SPEC, VARIANT, None, "square", tmp_path)


def test_render_image_background_adds_picture(tmp_path):
    img = tmp_path / "bg.png"
    img.write_bytes(base64.b64decode(PNG_1PX))
    spec = {**VALID_SPEC, "archetype": "hero_top", "background_style": "image"}
    p = render(spec, VARIANT, img, "portrait", tmp_path)
    prs, slide, _ = _texts(p)
    pics = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
    assert len(pics) >= 1


def test_render_image_style_without_image_still_renders(tmp_path):
    spec = {**VALID_SPEC, "background_style": "image"}
    p = render(spec, VARIANT, None, "portrait", tmp_path)
    prs, slide, joined = _texts(p)
    pics = [s for s in slide.shapes if s.shape_type == 13]
    assert len(pics) == 0
    assert VARIANT["headline"] in joined


def test_render_has_card_containers(tmp_path):
    spec = {**VALID_SPEC, "background_style": "solid", "card_style": "filled"}
    p = render(spec, VARIANT, None, "portrait", tmp_path)
    prs, slide, _ = _texts(p)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    autoshapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(autoshapes) >= len(VARIANT["points"])


def test_render_heading_font_applied(tmp_path):
    spec = {**VALID_SPEC, "background_style": "solid"}
    p = render(spec, VARIANT, None, "portrait", tmp_path)
    prs, slide, _ = _texts(p)
    fonts = set()
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts.add(run.font.name)
    assert "Anton" in fonts


def test_render_legacy_content_shape(tmp_path):
    legacy = {
        "headline": "H", "subheadline": "S", "facts": ["f1", "f2", "f3"], "cta": "C",
        "palette": {"bg": "#0E3A5D", "accent": "#3EC1D3", "text": "#FFFFFF"},
    }
    spec = {**VALID_SPEC, "background_style": "solid", "fact_stats": ["a", "b", "c"]}
    p = render(spec, legacy, None, "portrait", tmp_path)
    _, _, joined = _texts(p)
    for f in legacy["facts"]:
        assert f in joined
