import uuid
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

A4_EMU = {
    "portrait": (7_560_000, 10_692_000),
    "landscape": (10_692_000, 7_560_000),
}


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#"))


def _set_fill_alpha(shape, transparency_pct: int) -> None:
    """Add alpha to a solid fill. transparency_pct=40 → 60% opaque.

    Primary variant: shape.fill.fore_color._xFill is the <a:solidFill> element
    in python-pptx 1.0.2. If it fails, falls back to traversing _xPr.
    """
    try:
        srgb = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
        if srgb is None:
            raise AttributeError("srgbClr not found via _xFill")
    except AttributeError:
        # Fallback for python-pptx versions where _xFill differs
        srgb = shape.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    alpha = srgb.makeelement(qn("a:alpha"), {"val": str((100 - transparency_pct) * 1000)})
    srgb.append(alpha)


def _add_rect(slide, x, y, w, h, hex_color, transparency_pct=0):
    from pptx.enum.shapes import MSO_SHAPE
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = _rgb(hex_color)
    rect.line.fill.background()
    if transparency_pct:
        _set_fill_alpha(rect, transparency_pct)
    return rect


def _add_text(slide, x, y, w, h, text, size_pt, hex_color, bold=False, align=PP_ALIGN.LEFT, font=None, wrap=True, middle=False):
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    if middle:
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(hex_color)
    if font:
        run.font.name = font
    return box


_ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


def _add_auto_shape(slide, shape_name, x, y, w, h, hex_color, opacity=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    enum = {"rect": MSO_SHAPE.RECTANGLE, "line": MSO_SHAPE.RECTANGLE,
            "ellipse": MSO_SHAPE.OVAL}.get(shape_name, MSO_SHAPE.RECTANGLE)
    shp = slide.shapes.add_shape(enum, Emu(x), Emu(y), Emu(max(w, 1)), Emu(max(h, 1)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(hex_color)
    shp.line.fill.background()
    transparency_pct = int(round((1.0 - opacity) * 100))
    if transparency_pct:
        _set_fill_alpha(shp, transparency_pct)
    return shp


def _readable(hex_str: str) -> str:
    """Pick black or white text for legibility on the given background color."""
    hs = hex_str.lstrip("#")
    r, g, b = int(hs[0:2], 16), int(hs[2:4], 16), int(hs[4:6], 16)
    return "#12161C" if (0.299 * r + 0.587 * g + 0.114 * b) > 150 else "#FFFFFF"


def _fit(base_pt: int, text: str, comfortable_chars: int) -> int:
    """Shrink a font size for long text so it stays inside its box instead of
    overflowing (python-pptx does no auto-fit). Never below 60% of base."""
    n = len(text)
    if n <= comfortable_chars:
        return base_pt
    scale = max(0.6, comfortable_chars / n)
    return max(9, int(base_pt * scale))


def _fit_box(base_pt: int, text: str, w_emu: int, h_emu: int, min_pt: int = 10) -> int:
    """Size text to actually FILL a box: start at base_pt and step down only
    until the wrapped text fits the box's width AND height. Gives large,
    space-filling text in roomy cards instead of a timid fixed size."""
    import math
    w_pt = max(1, w_emu / 12700)
    h_pt = max(1, h_emu / 12700)
    n = max(1, len(text))
    for pt in range(base_pt, min_pt - 1, -1):
        chars_per_line = max(1, int(w_pt / (pt * 0.52)))
        lines = math.ceil(n / chars_per_line)
        if lines * pt * 1.22 <= h_pt:
            return pt
    return min_pt


def _rounded(slide, x, y, w, h, hex_fill, opacity=1.0, line_hex=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Emu(int(x)), Emu(int(y)), Emu(int(max(w, 1))), Emu(int(max(h, 1))))
    if opacity <= 0:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(hex_fill)
        if opacity < 1:
            _set_fill_alpha(shp, int(round((1 - opacity) * 100)))
    if line_hex:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Emu(19050)
    else:
        shp.line.fill.background()
    return shp


def _scrim(slide, x, y, w, h, hex_color, opacity):
    """Translucent overlay so text stays legible over imagery."""
    _add_rect(slide, int(x), int(y), int(w), int(h), hex_color,
              transparency_pct=int(round((1 - opacity) * 100)))


def _pic_cover(slide, image, x, y, w, h):
    slide.shapes.add_picture(str(image), Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))


def _draw_headline(slide, x, y, w, h, text, spec, color, size_pt, align="left", panel=False):
    pal, hs = spec["palette"], spec["header_style"]
    if panel:
        # solid band behind the headline so it never fights a busy image
        _add_rect(slide, int(x - w * 0.02), int(y - h * 0.12), int(w * 1.04), int(h * 1.3),
                  pal["bg"], transparency_pct=18)
    if hs == "underline":
        _add_rect(slide, int(x), int(y + h * 0.92), int(w * 0.32), max(int(h * 0.05), 30000), pal["accent"])
    elif hs == "badge":
        # small kicker bar ABOVE the headline — never a box behind the words
        _add_rect(slide, int(x), int(y - h * 0.14), int(w * 0.16), max(int(h * 0.05), 30000), pal["accent"])
    else:  # block (default): accent rule to the left of the headline
        _add_rect(slide, int(x), int(y + h * 0.05), int(w * 0.02), int(h * 0.8), pal["accent"])
        x = x + int(w * 0.045)
        w = w - int(w * 0.045)
    _add_text(slide, int(x), int(y), int(w), int(h), text, _fit(size_pt, text, 26), color,
              bold=True, align=_ALIGN_MAP.get(align, PP_ALIGN.LEFT), font=spec["fonts"]["heading"])


def _badge_w(badge_h, stat) -> int:
    """Chip width scaled to character count, so multi-char stats always get a
    single horizontal line (never a cramped vertical stack)."""
    n = len(str(stat))
    if n <= 1:
        return badge_h
    return min(int(badge_h * (1.0 + 0.52 * (n - 1))), int(badge_h * 2.7))


def _stat_badge(slide, x, y, badge_h, stat, accent, heading_font):
    """Draw a stat chip sized to its content; circle only for a single char,
    otherwise a rounded pill. Long stats are truncated and the font is scaled
    to the chip width so text is always a single, contained line."""
    stat = str(stat).strip()
    if len(stat) > 10:  # stats are meant to be punchy — clip label-like values
        stat = stat[:10].rstrip()
    bw = _badge_w(badge_h, stat)
    if len(stat) <= 1:
        _add_auto_shape(slide, "ellipse", x, y, bw, badge_h, accent)
    else:
        _rounded(slide, x, y, bw, badge_h, accent)
    # font must fit BOTH the chip height and its width
    by_height = int(badge_h / 12700 * 0.4)
    by_width = int(bw / 12700 / max(1, len(stat)) * 1.5)
    stat_pt = max(8, min(15, by_height, by_width))
    _add_text(slide, x, y, bw, badge_h, stat, stat_pt, _readable(accent),
              bold=True, align=PP_ALIGN.CENTER, font=heading_font, wrap=False, middle=True)
    return bw


def _card(slide, x, y, w, h, spec, stat, text, show_badge=True):
    pal, style = spec["palette"], spec["card_style"]
    surface, accent = pal["surface"], pal["accent"]
    pad = int(min(w, h) * 0.09)
    if style == "soft_shadow":
        off = int(min(w, h) * 0.035)
        _rounded(slide, x + off, y + off, w, h, pal["muted"], opacity=0.35)
        _rounded(slide, x, y, w, h, surface, opacity=1.0)
    elif style == "outline":
        _rounded(slide, x, y, w, h, surface, opacity=0.14, line_hex=accent)
    else:  # filled
        _rounded(slide, x, y, w, h, surface, opacity=1.0)
    text_color = _readable(surface)
    stat = str(stat)
    if not (show_badge and stat):
        tw, th = w - 2 * pad, h - pad
        _add_text(slide, x + pad, y + int(pad * 0.5), tw, th, text,
                  _fit_box(20, text, tw, th), text_color, align=PP_ALIGN.LEFT,
                  font=spec["fonts"]["body"], middle=True)
        return

    heading = spec["fonts"]["heading"]
    badge_h = min(int(h * 0.32), int(w * 0.16))
    badge_h = max(badge_h, 200000)
    if badge_h > int(h * 0.6):
        badge_h = int(h * 0.6)
    est_bw = _badge_w(badge_h, stat)

    if h >= badge_h * 2.4 and (w - est_bw - 3 * pad) < int(w * 0.55):
        # tall, narrowish card → stack the chip above the text
        bx, by = x + pad, y + pad
        _stat_badge(slide, bx, by, badge_h, stat, accent, heading)
        ty = by + badge_h + int(pad * 0.4)
        tw, th = w - 2 * pad, y + h - pad - ty
        _add_text(slide, x + pad, ty, tw, th, text,
                  _fit_box(19, text, tw, th), text_color, align=PP_ALIGN.LEFT,
                  font=spec["fonts"]["body"], middle=True)
    else:
        # side-by-side chip; text takes the rest, vertically centered
        bx, by = x + pad, y + (h - badge_h) // 2
        bw = _stat_badge(slide, bx, by, badge_h, stat, accent, heading)
        tx = bx + bw + pad
        tw, th = x + w - tx - pad, h - int(pad * 0.8)
        _add_text(slide, tx, y + int(pad * 0.4), tw, th, text,
                  _fit_box(18, text, tw, th), text_color, align=PP_ALIGN.LEFT,
                  font=spec["fonts"]["body"], middle=True)


def _draw_facts(slide, x, y, w, h, facts, stats, spec, cols):
    n = len(facts)
    if n == 0:
        return
    cols = max(1, min(cols, n))
    rows = (n + cols - 1) // cols
    gap = int(min(w, h) * 0.04)
    cw = (w - (cols - 1) * gap) // cols
    ch = (h - (rows - 1) * gap) // rows
    # cap card height so few-item layouts don't balloon into empty boxes;
    # then center the whole stack vertically in the region
    ch = min(ch, int(cw * 0.42) if cols == 1 else int(cw * 0.9))
    total = rows * ch + (rows - 1) * gap
    y0 = y + max(0, (h - total) // 2)
    for i, fact in enumerate(facts):
        r, c = divmod(i, cols)
        cx = x + c * (cw + gap)
        cy = y0 + r * (ch + gap)
        stat = stats[i] if i < len(stats) and str(stats[i]).strip() else str(i + 1)
        _card(slide, cx, cy, cw, ch, spec, stat, fact)


def _draw_cta(slide, x, y, w, h, text, spec):
    pal = spec["palette"]
    _rounded(slide, x, y, w, h, pal["accent"], opacity=1.0)
    _add_text(slide, int(x + w * 0.04), y, int(w * 0.92), h, text, _fit(22, text, 24),
              _readable(pal["accent"]), bold=True, align=PP_ALIGN.CENTER,
              font=spec["fonts"]["heading"], middle=True)


def _decor(slide, w, h, spec):
    """Fill non-image backgrounds with depth: a gradient-look wash plus accent
    blobs, so 'gradient'/'solid' posters stop rendering as one flat colour."""
    pal = spec["palette"]
    if spec["background_style"] == "gradient":
        _add_rect(slide, 0, 0, w, int(h * 0.55), pal["accent"], transparency_pct=78)
        _add_rect(slide, 0, int(h * 0.55), w, int(h * 0.2), pal["accent"], transparency_pct=90)
    if spec.get("accent_shapes", True):
        _add_auto_shape(slide, "ellipse", int(w * 0.72), int(-h * 0.06), int(w * 0.42), int(w * 0.42),
                        pal["accent"], opacity=0.16)
        _add_auto_shape(slide, "ellipse", int(-w * 0.12), int(h * 0.78), int(w * 0.34), int(w * 0.34),
                        pal["accent"], opacity=0.12)
    motif = spec.get("motif", "none")
    if motif == "dots":
        step = int(w * 0.045)
        for r in range(3):
            for c in range(6):
                _add_auto_shape(slide, "ellipse", int(w * 0.62) + c * step, int(h * 0.03) + r * step,
                                int(w * 0.008), int(w * 0.008), pal["accent"], opacity=0.55)
    elif motif == "rings":
        for rx, ry, rs in [(0.8, 0.05, 0.16), (0.86, 0.12, 0.1), (0.06, 0.86, 0.12)]:
            _add_auto_shape(slide, "ellipse", int(w * rx), int(h * ry), int(w * rs), int(w * rs),
                            pal["accent"], opacity=0.28)
            _add_auto_shape(slide, "ellipse", int(w * rx) + int(w * rs * 0.18), int(h * ry) + int(w * rs * 0.18),
                            int(w * rs * 0.64), int(w * rs * 0.64), pal["bg"], opacity=1.0)
    elif motif == "stripes":
        for i in range(4):
            _add_rect(slide, int(w * 0.7) + i * int(w * 0.035), int(-h * 0.02),
                      int(w * 0.012), int(h * 0.16), pal["accent"], transparency_pct=30)


def _draw_sources(slide, w, h, sources, spec, on_dark_bg):
    if not sources:
        return
    color = "#E8E8E8" if on_dark_bg else "#5A6472"
    _add_text(slide, int(w * 0.06), int(h * 0.965), int(w * 0.88), int(h * 0.03),
              "Sources: " + " · ".join(sources), 8, color, font=spec["fonts"]["body"])


def render(spec: dict, content: dict, image: Path | None, orientation: str, out_dir: Path) -> Path:
    """Render a StyleSpec (from app.design) + a content variant into an
    editable PPTX. Geometry is computed here per archetype, so output is
    always clean and full."""
    if orientation not in A4_EMU:
        raise ValueError(f"orientation must be one of {sorted(A4_EMU)}")
    w, h = A4_EMU[orientation]
    pal = spec["palette"]
    points = content.get("points")
    if points:
        facts = [p["text"] for p in points]
        stats = [p["stat"] for p in points]
    else:  # legacy content shape
        facts = content["facts"]
        stats = spec.get("fact_stats", [])
    sources = content.get("sources", [])
    arch = spec["archetype"]
    has_img = spec["background_style"] == "image" and image is not None
    m = int(w * 0.06)
    landscape = orientation == "landscape"

    prs = Presentation()
    prs.slide_width = Emu(w)
    prs.slide_height = Emu(h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_rect(slide, 0, 0, w, h, pal["bg"])  # base
    if not has_img:
        _decor(slide, w, h, spec)

    if arch == "hero_top":
        ih = int(h * 0.44)
        if has_img:
            _pic_cover(slide, image, 0, 0, w, ih)
            # lighter top scrim keeps the image visible; stronger only under text
            _scrim(slide, 0, int(ih * 0.55), w, int(ih * 0.45), pal["bg"], 0.5)
            head_color = "#FFFFFF"
        else:
            _add_rect(slide, 0, 0, w, ih, pal["accent"])
            head_color = _readable(pal["accent"])
        _draw_headline(slide, m, int(h * 0.13), w - 2 * m, int(h * 0.18), content["headline"],
                       spec, head_color, 42 if not landscape else 34, "left", panel=has_img)
        _add_text(slide, m, int(h * 0.47), w - 2 * m, int(h * 0.06), content["subheadline"],
                  _fit(18, content["subheadline"], 60), _readable(pal["bg"]), font=spec["fonts"]["body"])
        _draw_facts(slide, m, int(h * 0.55), w - 2 * m, int(h * 0.30), facts, stats, spec,
                    cols=2 if (landscape or len(facts) > 3) else 1)
        _draw_cta(slide, m, int(h * 0.88), w - 2 * m, int(h * 0.09), content["cta"], spec)

    elif arch == "sidebar":
        pw = int(w * 0.4)
        if has_img:
            _pic_cover(slide, image, 0, 0, pw, h)
            _scrim(slide, 0, 0, pw, h, pal["bg"], 0.6)
            head_color = "#FFFFFF"
        else:
            _add_rect(slide, 0, 0, pw, h, pal["accent"])
            head_color = _readable(pal["accent"])
        _draw_headline(slide, int(w * 0.05), int(h * 0.08), int(w * 0.3), int(h * 0.34),
                       content["headline"], spec, head_color, 34, "left")
        _draw_cta(slide, int(w * 0.05), int(h * 0.82), int(w * 0.3), int(h * 0.1), content["cta"], spec)
        _add_text(slide, int(w * 0.44), int(h * 0.08), int(w * 0.5), int(h * 0.1), content["subheadline"],
                  18, _readable(pal["bg"]), font=spec["fonts"]["body"])
        _draw_facts(slide, int(w * 0.44), int(h * 0.22), int(w * 0.5), int(h * 0.64), facts, stats, spec,
                    cols=2 if landscape else 1)

    elif arch == "banner_header":
        if has_img:
            _pic_cover(slide, image, 0, 0, w, h)
            _scrim(slide, 0, 0, w, h, pal["bg"], 0.68)
        bh = int(h * 0.2)
        _add_rect(slide, 0, 0, w, bh, pal["accent"])
        _draw_headline(slide, m, int(h * 0.03), w - 2 * m, int(h * 0.1), content["headline"],
                       spec, _readable(pal["accent"]), 38 if not landscape else 32, "center")
        _add_text(slide, m, int(h * 0.135), w - 2 * m, int(h * 0.05), content["subheadline"],
                  16, _readable(pal["accent"]), align=PP_ALIGN.CENTER, font=spec["fonts"]["body"])
        _draw_facts(slide, m, int(h * 0.25), w - 2 * m, int(h * 0.55), facts, stats, spec, cols=2)
        _draw_cta(slide, 0, int(h * 0.85), w, int(h * 0.15), content["cta"], spec)

    elif arch == "centered_feature":
        if has_img:
            _pic_cover(slide, image, 0, 0, w, h)
            _scrim(slide, 0, 0, w, h, pal["bg"], 0.68)
        elif spec["accent_shapes"]:
            _add_auto_shape(slide, "ellipse", int(w * 0.55), int(-h * 0.1), int(w * 0.6), int(w * 0.6),
                            pal["accent"], opacity=0.25)
        _draw_headline(slide, m, int(h * 0.26), w - 2 * m, int(h * 0.2), content["headline"],
                       spec, "#FFFFFF" if has_img else _readable(pal["bg"]),
                       50 if not landscape else 40, "center", panel=has_img)
        _add_text(slide, m, int(h * 0.46), w - 2 * m, int(h * 0.07), content["subheadline"],
                  18, "#FFFFFF" if has_img else _readable(pal["bg"]), align=PP_ALIGN.CENTER,
                  font=spec["fonts"]["body"])
        _draw_facts(slide, m, int(h * 0.55), w - 2 * m, int(h * 0.30), facts, stats, spec,
                    cols=2 if len(facts) > 2 else 1)
        _draw_cta(slide, int(w * 0.25), int(h * 0.86), int(w * 0.5), int(h * 0.09), content["cta"], spec)

    elif arch == "big_number":
        # lead stat becomes the poster's hero element
        if has_img:
            _pic_cover(slide, image, 0, 0, w, h)
            _scrim(slide, 0, 0, w, h, pal["bg"], 0.68)
        base_color = "#FFFFFF" if has_img else _readable(pal["bg"])
        # hero stat: a single short token only, never a phrase that could wrap
        lead = str(stats[0]).split()[0] if stats and str(stats[0]).strip() else "!"
        lead = lead[:8]
        hero_pt = (110 if len(lead) <= 5 else 78) if not landscape else 84
        _add_text(slide, m, int(h * 0.05), w - 2 * m, int(h * 0.19), lead, hero_pt,
                  pal["accent"], bold=True, align=PP_ALIGN.LEFT, font=spec["fonts"]["heading"],
                  wrap=False, middle=True)
        _draw_headline(slide, m, int(h * 0.27), w - 2 * m, int(h * 0.12), content["headline"],
                       spec, base_color, 38 if not landscape else 32, "left")
        _add_text(slide, m, int(h * 0.39), w - 2 * m, int(h * 0.05), content["subheadline"],
                  16, base_color, font=spec["fonts"]["body"])
        # the hero stat's own point text stays visible as the lead line
        _add_text(slide, m, int(h * 0.445), w - 2 * m, int(h * 0.05), facts[0],
                  15, pal["accent"] if not has_img else "#FFFFFF", bold=True,
                  font=spec["fonts"]["body"])
        rest = list(zip(facts, stats))[1:] if len(facts) > 1 else list(zip(facts, stats))
        _draw_facts(slide, m, int(h * 0.51), w - 2 * m, int(h * 0.33),
                    [f for f, _ in rest], [s for _, s in rest], spec,
                    cols=2 if (landscape or len(rest) > 3) else 1)
        _draw_cta(slide, m, int(h * 0.87), w - 2 * m, int(h * 0.09), content["cta"], spec)

    elif arch == "steps_path":
        # numbered path down the page with a connecting spine — fits checklists/steps
        if has_img:
            _pic_cover(slide, image, 0, 0, w, h)
            _scrim(slide, 0, 0, w, h, pal["bg"], 0.72)
        base_color = "#FFFFFF" if has_img else _readable(pal["bg"])
        _draw_headline(slide, m, int(h * 0.05), w - 2 * m, int(h * 0.1), content["headline"],
                       spec, base_color, 36 if not landscape else 30, "left")
        _add_text(slide, m, int(h * 0.155), w - 2 * m, int(h * 0.05), content["subheadline"],
                  15, base_color, font=spec["fonts"]["body"])
        top, bottom = int(h * 0.24), int(h * 0.84)
        n = len(facts)
        spine_x = m + int(w * 0.045)
        _add_rect(slide, spine_x - 15000, top, 30000, bottom - top - int((bottom - top) / max(n, 1) * 0.35),
                  pal["accent"])
        row = (bottom - top) // max(n, 1)
        badge = int(min(row * 0.42, w * 0.09))
        card_h = min(int(row * 0.62), int((w - m - spine_x - badge) * 0.26))
        for i, fact in enumerate(facts):
            cy = top + i * row + (row - card_h) // 2
            node_y = cy + (card_h - badge) // 2
            _add_auto_shape(slide, "ellipse", spine_x - badge // 2, node_y, badge, badge, pal["accent"])
            _add_text(slide, spine_x - badge // 2, node_y, badge, badge, str(i + 1),
                      16, _readable(pal["accent"]), bold=True, align=PP_ALIGN.CENTER,
                      font=spec["fonts"]["heading"], wrap=False, middle=True)
            card_x = spine_x + badge
            # node numbers the step; card shows the stat as a small inline badge + text
            _card(slide, card_x, cy, w - m - card_x, card_h, spec,
                  stats[i] if i < len(stats) else "", fact,
                  show_badge=bool(stats[i]) if i < len(stats) else False)
        _draw_cta(slide, m, int(h * 0.88), w - 2 * m, int(h * 0.08), content["cta"], spec)

    elif arch == "poster_frame":
        # thick brand-color frame around the whole poster, gallery style
        bw = int(min(w, h) * 0.035)
        if has_img:
            _pic_cover(slide, image, 0, 0, w, h)
            _scrim(slide, 0, 0, w, h, pal["bg"], 0.66)
        for fx, fy, fw2, fh2 in [(0, 0, w, bw), (0, h - bw, w, bw), (0, 0, bw, h), (w - bw, 0, bw, h)]:
            _add_rect(slide, fx, fy, fw2, fh2, pal["accent"])
        base_color = "#FFFFFF" if has_img else _readable(pal["bg"])
        inner = bw + m
        _draw_headline(slide, inner, int(h * 0.09), w - 2 * inner, int(h * 0.13), content["headline"],
                       spec, base_color, 40 if not landscape else 33, "center", panel=has_img)
        _add_text(slide, inner, int(h * 0.24), w - 2 * inner, int(h * 0.05), content["subheadline"],
                  _fit(16, content["subheadline"], 60), base_color, align=PP_ALIGN.CENTER,
                  font=spec["fonts"]["body"])
        _draw_facts(slide, inner, int(h * 0.32), w - 2 * inner, int(h * 0.5), facts, stats, spec,
                    cols=2 if (landscape or len(facts) > 4) else 1)
        _draw_cta(slide, int(w * 0.22), int(h * 0.855), int(w * 0.56), int(h * 0.08), content["cta"], spec)

    elif arch == "diagonal_energy":
        # angled accent slabs give the poster motion; content rides the diagonal
        if has_img:
            _pic_cover(slide, image, 0, 0, w, h)
            _scrim(slide, 0, 0, w, h, pal["bg"], 0.68)
        from pptx.enum.shapes import MSO_SHAPE
        tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Emu(0), Emu(0), Emu(w), Emu(int(h * 0.34)))
        tri.rotation = 180
        tri.fill.solid()
        tri.fill.fore_color.rgb = _rgb(pal["accent"])
        tri.line.fill.background()
        _set_fill_alpha(tri, 22)
        tri2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Emu(0), Emu(int(h * 0.72)),
                                      Emu(w), Emu(int(h * 0.28)))
        tri2.fill.solid()
        tri2.fill.fore_color.rgb = _rgb(pal["accent"])
        tri2.line.fill.background()
        _set_fill_alpha(tri2, 35)
        base_color = "#FFFFFF" if has_img else _readable(pal["bg"])
        _draw_headline(slide, m, int(h * 0.06), w - 2 * m, int(h * 0.14), content["headline"],
                       spec, base_color, 42 if not landscape else 34, "left", panel=has_img)
        _add_text(slide, m, int(h * 0.215), w - 2 * m, int(h * 0.05), content["subheadline"],
                  _fit(16, content["subheadline"], 60), base_color, font=spec["fonts"]["body"])
        _draw_facts(slide, m, int(h * 0.29), w - 2 * m, int(h * 0.5), facts, stats, spec,
                    cols=2 if (landscape or len(facts) > 3) else 1)
        _draw_cta(slide, int(w * 0.3), int(h * 0.86), int(w * 0.64), int(h * 0.09), content["cta"], spec)

    elif arch == "bottom_hero":
        # inverse hero: content on top, imagery anchors the bottom third
        ih = int(h * 0.34)
        if has_img:
            _pic_cover(slide, image, 0, h - ih, w, ih)
            _scrim(slide, 0, h - ih, w, int(ih * 0.4), pal["bg"], 0.5)
        else:
            _add_rect(slide, 0, h - ih, w, ih, pal["accent"])
        base_color = _readable(pal["bg"])
        _draw_headline(slide, m, int(h * 0.05), w - 2 * m, int(h * 0.13), content["headline"],
                       spec, base_color, 42 if not landscape else 34, "left")
        _add_text(slide, m, int(h * 0.19), w - 2 * m, int(h * 0.05), content["subheadline"],
                  _fit(16, content["subheadline"], 60), base_color, font=spec["fonts"]["body"])
        _draw_facts(slide, m, int(h * 0.26), w - 2 * m, int(h * 0.36), facts, stats, spec,
                    cols=2 if (landscape or len(facts) > 4) else 1)
        _draw_cta(slide, m, int(h - ih + ih * 0.32), w - 2 * m, int(ih * 0.36), content["cta"], spec)

    else:  # split_band — full-height image band beside a content column
        band = int(w * 0.46)
        cx0 = band + m // 2  # content column start
        if has_img:
            _pic_cover(slide, image, 0, 0, band, h)
            _scrim(slide, 0, int(h * 0.7), band, int(h * 0.3), pal["bg"], 0.6)
        else:
            _add_rect(slide, 0, 0, band, h, pal["accent"])
            # fill the empty solid panel with a large reversed stat + motif so
            # it never reads as dead space
            _add_auto_shape(slide, "ellipse", int(band * 0.15), int(h * 0.5), int(band * 0.8),
                            int(band * 0.8), pal["bg"], opacity=0.18)
            lead = stats[0] if stats else ""
            if lead:
                _add_text(slide, int(band * 0.08), int(h * 0.16), int(band * 0.84), int(h * 0.22),
                          str(lead)[:6], _fit(64, str(lead)[:6], 5), _readable(pal["accent"]),
                          bold=True, align=PP_ALIGN.CENTER, font=spec["fonts"]["heading"], wrap=False)
            _add_text(slide, int(band * 0.1), int(h * 0.42), int(band * 0.8), int(h * 0.3),
                      content["subheadline"], _fit(16, content["subheadline"], 50),
                      _readable(pal["accent"]), align=PP_ALIGN.CENTER, font=spec["fonts"]["body"], middle=True)
        base_color = _readable(pal["bg"])
        _draw_headline(slide, cx0, int(h * 0.07), w - cx0 - m, int(h * 0.16), content["headline"],
                       spec, base_color, 34 if not landscape else 30, "left")
        _add_text(slide, cx0, int(h * 0.25), w - cx0 - m, int(h * 0.06), content["subheadline"],
                  15, base_color, font=spec["fonts"]["body"])
        _draw_facts(slide, cx0, int(h * 0.33), w - cx0 - m, int(h * 0.5), facts, stats, spec, cols=1)
        _draw_cta(slide, cx0, int(h * 0.86), w - cx0 - m, int(h * 0.09), content["cta"], spec)

    _draw_sources(slide, w, h, sources, spec, on_dark_bg=(has_img or _readable(pal["bg"]) == "#FFFFFF"))

    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / f"poster_{uuid.uuid4().hex[:8]}.pptx"
    prs.save(str(path))
    return path


def build(content: dict, image: Path | None, orientation: str, out_dir: Path) -> Path:
    if orientation not in A4_EMU:
        raise ValueError(f"orientation must be one of {sorted(A4_EMU)}")
    w, h = A4_EMU[orientation]
    pal = content["palette"]

    prs = Presentation()
    prs.slide_width = Emu(w)
    prs.slide_height = Emu(h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    if image is not None:
        slide.shapes.add_picture(str(image), 0, 0, Emu(w), Emu(h))
        _add_rect(slide, 0, 0, w, h, pal["bg"], transparency_pct=55)
    else:
        _add_rect(slide, 0, 0, w, h, pal["bg"])
        _add_rect(slide, 0, int(h * 0.92), w, int(h * 0.08), pal["accent"])
        _add_rect(slide, int(w * 0.82), 0, int(w * 0.18), int(h * 0.25), pal["accent"], transparency_pct=30)

    margin = int(w * 0.08)
    cw = w - 2 * margin
    headline_size = 44 if orientation == "portrait" else 40

    _add_text(slide, margin, int(h * 0.07), cw, int(h * 0.14),
              content["headline"], headline_size, pal["text"], bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, margin, int(h * 0.21), cw, int(h * 0.08),
              content["subheadline"], 20, pal["text"], align=PP_ALIGN.CENTER)

    facts_top = int(h * 0.34)
    row_h = int(h * 0.40 / len(content["facts"]))
    for i, fact in enumerate(content["facts"]):
        _add_rect(slide, margin, facts_top + i * row_h, int(w * 0.012), int(row_h * 0.72), pal["accent"])
        _add_text(slide, margin + int(w * 0.03), facts_top + i * row_h, cw - int(w * 0.03), int(row_h * 0.8),
                  fact, 16, pal["text"])

    cta_top = int(h * 0.80)
    _add_rect(slide, margin, cta_top, cw, int(h * 0.10), pal["accent"])
    _add_text(slide, margin, cta_top + int(h * 0.025), cw, int(h * 0.05),
              content["cta"], 24, pal["bg"], bold=True, align=PP_ALIGN.CENTER)

    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / f"poster_{uuid.uuid4().hex[:8]}.pptx"
    prs.save(str(path))
    return path


# Kept as the safety net when AI layout generation or validation fails.
fallback_build = build
